#!/usr/bin/env python3
"""
Create a 2-page US Election themed PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ============================================
    # Slide 1: Title Slide
    # ============================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Red header bar
    shape1 = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(197, 48, 48)  # Deep red
    shape1.line.fill.background()
    
    # Title text
    title_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "2024 美国大选"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), prs.slide_width - Inches(2), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "总统选举与政治格局分析"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(40, 50, 70)
    p.alignment = PP_ALIGN.CENTER
    
    # Year badge
    badge = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        prs.slide_width/2 - Inches(1.5), Inches(4), Inches(3), Inches(0.8)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
    badge.line.color.rgb = RGBColor(26, 54, 93)
    badge.line.width = Pt(2)
    
    badge_frame = badge.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = badge_frame.paragraphs[0]
    p.text = "2024"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 54, 93)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide1.shapes.add_textbox(
        Inches(1), Inches(6.5), prs.slide_width - Inches(2), Inches(0.5)
    )
    footer_frame = footer_box.text_frame
    footer_frame.word_wrap = True
    p = footer_frame.paragraphs[0]
    p.text = "民主党的挑战与共和党的策略"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================
    # Slide 2: Electoral System & Key Issues
    # ============================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Red accent bar at top
    accent = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(197, 48, 48)
    accent.line.fill.background()
    
    # Title
    title_box2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.7)
    )
    title_frame2 = title_box2.text_frame
    title_frame2.word_wrap = True
    p = title_frame2.paragraphs[0]
    p.text = "选举制度与关键议题"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 54, 93)
    
    # Page number indicator
    p2 = title_frame2.add_paragraph()
    p2.text = "第二页"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(197, 48, 48)
    p2.alignment = PP_ALIGN.RIGHT
    
    # Two column layout
    # Column 1: Electoral College
    col1_x = Inches(0.5)
    col2_x = Inches(5)
    col_width = Inches(6)
    content_y = Inches(1.5)
    
    # Section 1 header
    sec1_header = slide2.shapes.add_textbox(col1_x, content_y, col_width, Inches(0.5))
    sec1_frame = sec1_header.text_frame
    sec1_frame.word_wrap = True
    p = sec1_frame.paragraphs[0]
    p.text = "🗳️ 选举人团制度"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(45, 55, 72)
    # Underline
    line = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, col1_x, content_y + Inches(0.45), Inches(2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(197, 48, 48)
    line.line.fill.background()
    
    # Section 1 content
    content1_box = slide2.shapes.add_textbox(
        col1_x, content_y + Inches(0.7), col_width, Inches(3.5)
    )
    content1_frame = content1_box.text_frame
    content1_frame.word_wrap = True
    items1 = [
        "538张选举人票",
        "270票获胜门槛",
        "7个摇摆州决定胜负",
        "宾夕法尼亚、密歇根、威斯康星",
        "佐治亚、亚利桑那、内华达、北卡罗来纳"
    ]
    for i, item in enumerate(items1):
        if i == 0:
            p = content1_frame.paragraphs[0]
        else:
            p = content1_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(74, 85, 104)
        p.space_before = Pt(8)
    
    # Column 2: Key Issues
    sec2_header = slide2.shapes.add_textbox(col2_x, content_y, col_width, Inches(0.5))
    sec2_frame = sec2_header.text_frame
    sec2_frame.word_wrap = True
    p = sec2_frame.paragraphs[0]
    p.text = "📊 核心议题"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(45, 55, 72)
    # Underline
    line2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, col2_x, content_y + Inches(0.45), Inches(1.5), Inches(0.08)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(197, 48, 48)
    line2.line.fill.background()
    
    # Section 2 content
    content2_box = slide2.shapes.add_textbox(
        col2_x, content_y + Inches(0.7), col_width, Inches(3.5)
    )
    content2_frame = content2_box.text_frame
    content2_frame.word_wrap = True
    items2 = [
        "经济与通胀",
        "移民政策",
        "堕胎权争议",
        "外交与国家安全",
        "气候变化政策"
    ]
    for i, item in enumerate(items2):
        if i == 0:
            p = content2_frame.paragraphs[0]
        else:
            p = content2_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(74, 85, 104)
        p.space_before = Pt(8)
    
    # Highlight box
    highlight = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        col2_x, content_y + Inches(3), Inches(5.5), Inches(0.8)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(235, 248, 255)
    highlight.line.color.rgb = RGBColor(49, 130, 206)
    highlight.line.width = Pt(1)
    
    hl_frame = highlight.text_frame
    hl_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = hl_frame.paragraphs[0]
    p.text = "关键摇摆选民群体"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(43, 108, 176)
    p.alignment = PP_ALIGN.CENTER
    
    # Statistics row
    stats_y = Inches(5.8)
    stats_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(0.5), stats_y, prs.slide_width - Inches(1), Inches(1.2)
    )
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
    stats_box.line.color.rgb = RGBColor(197, 48, 48)
    stats_box.line.width = Pt(1)
    
    stats_frame = stats_box.text_frame
    stats_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = stats_frame.paragraphs[0]
    p.text = "2.4亿合格选民" + " " * 25 + "50州同步投票" + " " * 25 + "11月5日投票日"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 54, 93)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = '/Users/admin/work/agent_loop/us_election_2pages.pptx'
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
